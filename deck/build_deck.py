#!/usr/bin/env python3
"""Build the 'should i cop this?' pitch deck (Tangerine-orange theme) with live
app screenshots — including each persona's financial dashboard. Run:
   python deck/build_deck.py"""

from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
SHOTS = HERE / "shots"

CREAM = RGBColor(0xFD, 0xF4, 0xE9)
ORANGE = RGBColor(0xFF, 0x6A, 0x00)
ORANGE_DEEP = RGBColor(0xD9, 0x4E, 0x00)
INK = RGBColor(0x24, 0x17, 0x10)
SOFT = RGBColor(0x7A, 0x6A, 0x5B)
GREEN = RGBColor(0x15, 0x93, 0x5A)
YELLOW = RGBColor(0xC9, 0x8A, 0x00)
RED = RGBColor(0xDB, 0x3B, 0x3B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DISPLAY = "Bricolage Grotesque"
BODY = "Segoe UI"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def slide(bg=CREAM):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def text(s, txt, left, top, width, height, size, color=INK, bold=True,
         font=DISPLAY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line=1.0):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(txt.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = color
    return tb


def dot(s, left, top, size, color=ORANGE):
    sh = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()


def bar(s, left, top, width, height, color=ORANGE):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()


def pill(s, label, left, top, color):
    w = 0.55 + 0.16 * len(label)
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(w), Inches(0.6))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    p = sh.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = label; r.font.size = Pt(20); r.font.bold = True
    r.font.name = DISPLAY; r.font.color.rgb = WHITE
    return left + w + 0.25


def img(s, name, left, top, box_w, box_h, top_align=False, border=True):
    path = SHOTS / name
    iw, ih = Image.open(path).size
    ar = iw / ih
    w = min(box_w, box_h * ar)
    h = w / ar
    l = left + (box_w - w) / 2
    t = top if top_align else top + (box_h - h) / 2
    pic = s.shapes.add_picture(str(path), Inches(l), Inches(t), width=Inches(w))
    if border:
        pic.line.color.rgb = RGBColor(0xEF, 0xE2, 0xD2); pic.line.width = Pt(1)
    return l, t, w, h


# ---- S1 title ----
s = slide(CREAM)
bar(s, 0, 0, 13.333, 0.28, ORANGE)
dot(s, 0.9, 2.45, 0.7)
text(s, "should i cop this?", 1.8, 2.0, 11, 2.0, 78, INK, line=0.95)
text(s, "a gen z money agent that actually knows your money", 1.85, 3.7, 11, 0.8, 26, ORANGE_DEEP)
text(s, "ask by voice or text. get a real verdict — grounded in your actual account data.",
     1.85, 4.4, 11, 0.8, 18, SOFT, bold=False, font=BODY)
text(s, "TechTO · May 24, 2026      Chip · Luca · Abdul · Pablo", 1.85, 6.6, 11, 0.5, 15, SOFT, bold=False, font=BODY)

# ---- S2 problem ----
s = slide(CREAM)
text(s, "the problem", 0.9, 0.7, 11, 0.8, 22, ORANGE_DEEP)
text(s, "everyone asks “should i buy this?”\nnobody can answer in the moment.", 0.9, 1.5, 11.5, 2.2, 44, INK)
for i, it in enumerate([
    "money is the #1 stress for people — yet the math is invisible at checkout.",
    "budget apps show pie charts of last month. not a decision right now.",
    "“will this set me back?” takes a spreadsheet nobody opens.",
]):
    y = 4.0 + i * 0.75
    dot(s, 0.95, y + 0.08, 0.18)
    text(s, it, 1.35, y, 11, 0.7, 19, INK, bold=False, font=BODY)

# ---- S3 solution ----
s = slide(CREAM)
text(s, "what we built", 0.9, 0.7, 11, 0.8, 22, ORANGE_DEEP)
text(s, "a money bestie that gives you a straight verdict —\ngrounded in your real numbers, in a voice that slaps.",
     0.9, 1.5, 11.7, 2.0, 34, INK, line=1.05)
text(s, "one question. four answers:", 0.9, 3.7, 11, 0.6, 20, SOFT, bold=False, font=BODY)
x = 0.9
for lbl, col in [("cop it", GREEN), ("wait", YELLOW), ("skip", RED), ("drop", RED)]:
    x = pill(s, lbl, x, 4.4, col)
text(s, "voice or text · 99% deterministic engine, 1% AI · never invents a number",
     0.9, 5.6, 11.5, 0.6, 18, ORANGE_DEEP, font=BODY)

# ---- S4 the money shot ----
s = slide(INK)
text(s, "real data in. honest answer out. in seconds.", 0.7, 0.45, 12, 0.7, 24, WHITE)
img(s, "full-luca.png", 0.7, 1.25, 12, 6.0)

# ---- S5 four people, four wallets (the financial dashboards) ----
s = slide(CREAM)
text(s, "four people. four wallets. one agent.", 0.7, 0.5, 12, 0.7, 30, INK)
text(s, "every verdict reads their real Tangerine-style account — income, subs, goals, spending.",
     0.7, 1.2, 12, 0.5, 16, SOFT, bold=False, font=BODY)
people = [
    ("profile-luca.png", "luca", "19, barista · −$470/mo"),
    ("profile-daniel.png", "daniel", "26, dev · +$275/mo, saving for a home"),
    ("profile-chen.png", "the chens", "family of 4 · +$6.9k/mo"),
    ("profile-margaret.png", "margaret", "retired · +$1.3k/mo, owns home"),
]
col_w = 3.15
gap = 0.12
start = (13.333 - (4 * col_w + 3 * gap)) / 2
for i, (im, name, sub) in enumerate(people):
    left = start + i * (col_w + gap)
    text(s, name, left, 1.95, col_w, 0.4, 18, ORANGE_DEEP, align=PP_ALIGN.CENTER)
    text(s, sub, left, 2.35, col_w, 0.5, 12, SOFT, bold=False, font=BODY, align=PP_ALIGN.CENTER)
    img(s, im, left, 2.95, col_w, 4.35, top_align=True)

# ---- S6 same agent, different call ----
s = slide(CREAM)
text(s, "same agent. different person. different call.", 0.7, 0.5, 12.2, 0.7, 30, INK)
pairs = [("full-luca.png", "luca (broke student)", "skip", RED),
         ("full-chen.png", "the chens (comfortable family)", "cop it", GREEN)]
for i, (im, who, verdict, col) in enumerate(pairs):
    left = 0.6 + i * 6.25
    text(s, who, left, 1.35, 6.0, 0.4, 17, INK, align=PP_ALIGN.CENTER)
    img(s, im, left, 1.85, 6.0, 4.6)
    bar(s, left + 3.0 - 0.85, 6.55, 1.7, 0.55, col)
    text(s, verdict, left + 3.0 - 0.85, 6.62, 1.7, 0.4, 20, WHITE, align=PP_ALIGN.CENTER)

# ---- S7 how it works ----
s = slide(CREAM)
text(s, "how it works", 0.9, 0.7, 11, 0.8, 22, ORANGE_DEEP)
text(s, "“AI is the 1%. the engine is the 99%.”", 0.9, 1.45, 11.5, 0.9, 34, INK)
for i, (h, b) in enumerate([
    ("the engine", "deterministic rules over your real transactions: budget pace, dead subs, spending leaks, goal impact → the verdict. python, no guessing."),
    ("claude", "turns the verdict into a gen z answer. reads your real numbers, never invents one. cop / wait / skip / drop, in character."),
    ("voice", "talk to it out loud (ElevenLabs). it knows your money before you finish the sentence — and shows the card."),
]):
    left = 0.9 + i * 4.0
    bar(s, left, 3.0, 3.5, 0.08, ORANGE)
    text(s, h, left, 3.2, 3.6, 0.6, 22, ORANGE_DEEP)
    text(s, b, left, 3.95, 3.6, 2.6, 15, INK, bold=False, font=BODY, line=1.1)

# ---- S8 closing ----
s = slide(ORANGE)
dot(s, 0.9, 2.3, 0.6, WHITE)
text(s, "should i cop this?", 1.7, 1.95, 11, 1.4, 60, WHITE, line=0.95)
text(s, "a real-world problem everyone has weekly · a money app with an actual personality ·\na thoughtful agent that's grounded, not a chatbot.",
     1.7, 3.6, 11, 1.2, 19, WHITE, bold=False, font=BODY, line=1.2)
text(s, "Chip · Luca · Abdul · Pablo", 1.7, 6.4, 11, 0.5, 16, RGBColor(0xFF, 0xE3, 0xCB), font=BODY)

out = HERE / "should-i-cop-this.pptx"
prs.save(str(out))
print("saved", out, "·", len(prs.slides._sldIdLst), "slides")
